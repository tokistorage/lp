#!/usr/bin/env python3
"""
TokiStorage Client Proposal Deck Generator (Timeless Consulting)
Generates PPTX files (JP + EN) and converts to PDF via LibreOffice.
Design aligned with TokiStorage landing page (index.css).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import subprocess, os, sys, re, zipfile, io, tempfile

# ── Circular photo helper ─────────────────────────────────────────────

def _prepare_circular_photo(src_path, size=300):
    from PIL import Image, ImageDraw
    img = Image.open(src_path).convert("RGBA")
    w, h = img.size
    s = min(w, h)
    left, top = (w - s) // 2, (h - s) // 2
    img = img.crop((left, top, left + s, top + s)).resize((size, size), Image.LANCZOS)
    mask = Image.new("L", (size, size), 0)
    ImageDraw.Draw(mask).ellipse((0, 0, size, size), fill=255)
    result = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    result.paste(img, mask=mask)
    tmp = os.path.join(tempfile.gettempdir(), "_client_deck_profile.png")
    result.save(tmp, "PNG")
    return tmp

# ── Design tokens (matched to index.css :root) ────────────────────────
TOKI_BLUE      = RGBColor(0x25, 0x63, 0xEB)
TOKI_BLUE_DK   = RGBColor(0x1D, 0x4E, 0xD8)
TOKI_BLUE_PALE = RGBColor(0xEF, 0xF6, 0xFF)
GOLD           = RGBColor(0xC9, 0xA9, 0x62)
GOLD_PALE      = RGBColor(0xFD, 0xF8, 0xE8)
EMERALD        = RGBColor(0x16, 0xA3, 0x4A)
GREEN_PALE     = RGBColor(0xF0, 0xFD, 0xF4)

TEXT_PRIMARY   = RGBColor(0x1E, 0x29, 0x3B)
TEXT_SECONDARY = RGBColor(0x47, 0x55, 0x69)
TEXT_MUTED     = RGBColor(0x94, 0xA3, 0xB8)
WHITE          = RGBColor(0xFF, 0xFF, 0xFF)

BG_PAGE        = RGBColor(0xF8, 0xFA, 0xFC)
BG_SECTION     = RGBColor(0xF1, 0xF5, 0xF9)
BORDER         = RGBColor(0xE2, 0xE8, 0xF0)

DARK_BG        = RGBColor(0x1E, 0x29, 0x3B)
DARK_BG2       = RGBColor(0x0F, 0x17, 0x2A)

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)

FONT_JP = "IPAPGothic"
FONT_EN = "Calibri"

OUT_DIR = os.path.dirname(os.path.abspath(__file__))


# ── Helpers ────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs

def add_blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def add_rect(slide, left, top, width, height, fill=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(0.75)
    else:
        shape.line.fill.background()
    return shape

def set_text(tf, text, font_name, size, color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, line_spacing=1.35):
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def add_para(tf, text, font_name, size, color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, space_before=0, line_spacing=1.35):
    p = tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line_spacing
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p

def add_textbox(slide, left, top, width, height, text, font_name, size,
                color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        txBox.text_frame.paragraphs[0].space_before = Pt(0)
        txBox.text_frame.paragraphs[0].space_after = Pt(0)
    except:
        pass
    from pptx.oxml.ns import qn
    txBody = txBox.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if anchor == MSO_ANCHOR.MIDDLE:
        bodyPr.set('anchor', 'ctr')
    elif anchor == MSO_ANCHOR.BOTTOM:
        bodyPr.set('anchor', 'b')
    set_text(tf, text, font_name, size, color, bold, align)
    return txBox

def add_action_bar(slide, text, font):
    bar_h = Inches(0.65)
    add_rect(slide, 0, 0, SLIDE_W, bar_h, fill=DARK_BG)
    add_textbox(slide, Inches(0.5), 0, SLIDE_W - Inches(1), bar_h,
                text, font, 12, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)

def add_footer(slide, left_text, pg, font):
    y = SLIDE_H - Inches(0.38)
    add_rect(slide, 0, y, SLIDE_W, Pt(0.5), fill=BORDER)
    add_textbox(slide, Inches(0.5), y + Pt(2), Inches(4), Inches(0.3),
                left_text, font, 9, TEXT_MUTED)
    add_textbox(slide, Inches(4), y + Pt(2), Inches(2), Inches(0.3),
                "Confidential", font, 9, TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, SLIDE_W - Inches(1), y + Pt(2), Inches(0.5), Inches(0.3),
                str(pg), font, 9, TEXT_MUTED, bold=True, align=PP_ALIGN.RIGHT)

def add_section_label(slide, text, font, top):
    add_textbox(slide, Inches(0.5), top, Inches(3), Inches(0.3),
                text.upper(), font, 10, TOKI_BLUE, bold=True)


# ── Card helpers ──────────────────────────────────────────────────────

def draw_col_card(slide, x, y, w, h, num, title, body, font):
    add_rect(slide, x, y, w, h, fill=WHITE, border_color=BORDER)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(0.6), Inches(0.35),
                num, font, 15, TOKI_BLUE, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.45), w - Inches(0.3), Inches(0.75),
                title, font, 11, TEXT_PRIMARY, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(1.3), w - Inches(0.3), h - Inches(1.4),
                body, font, 11, TEXT_SECONDARY)

def draw_grid_card(slide, x, y, w, h, num, title, body, font):
    """Numbered card for offering section"""
    add_rect(slide, x, y, w, h, fill=WHITE, border_color=BORDER)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.08), Inches(0.6), Inches(0.3),
                num, font, 16, TOKI_BLUE, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.38), w - Inches(0.3), Inches(0.35),
                title, font, 11, TEXT_PRIMARY, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.78), w - Inches(0.3), h - Inches(0.88),
                body, font, 10, TEXT_SECONDARY)

def draw_sector_card(slide, x, y, w, h, icon_letter, title, body, font):
    add_rect(slide, x, y, w, h, fill=WHITE, border_color=BORDER)
    circ_x = x + (w - Inches(0.42)) / 2
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, circ_x, y + Inches(0.12), Inches(0.42), Inches(0.42))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TOKI_BLUE_PALE
    circ.line.fill.background()
    add_textbox(slide, circ_x, y + Inches(0.12), Inches(0.42), Inches(0.42),
                icon_letter, font, 12, TOKI_BLUE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, x + Inches(0.08), y + Inches(0.6), w - Inches(0.16), Inches(0.26),
                title, font, 11, TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.08), y + Inches(0.95), w - Inches(0.16), h - Inches(1.05),
                body, font, 9, TEXT_SECONDARY, align=PP_ALIGN.CENTER)

def draw_deliverable_item(slide, x, y, w, h, icon_bg, icon_letter, title, body, font):
    add_rect(slide, x, y, w, h, fill=WHITE, border_color=BORDER)
    # Icon square
    ix, iy = x + Inches(0.15), y + Inches(0.15)
    icon_box = add_rect(slide, ix, iy, Inches(0.5), Inches(0.5), fill=icon_bg)
    add_textbox(slide, ix, iy, Inches(0.5), Inches(0.5),
                icon_letter, font, 14, TOKI_BLUE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, x + Inches(0.8), y + Inches(0.12), w - Inches(1.0), Inches(0.32),
                title, font, 11, TEXT_PRIMARY, bold=True)
    add_textbox(slide, x + Inches(0.8), y + Inches(0.46), w - Inches(1.0), h - Inches(0.58),
                body, font, 10, TEXT_SECONDARY)


# ══════════════════════════════════════════════════════════════════════
#  CONTENT DATA
# ══════════════════════════════════════════════════════════════════════

CONTENT = {
    "ja": {
        "font": FONT_JP,
        "filename": "tokistorage-client-deck",
        "cover": {
            "label": "Timeless Consulting",
            "title": "100年後に、\n何を残しますか。",
            "sub": "AIがあらゆる問いに答える時代。\nしかしこの問いだけは、あなた自身にしか答えられない。",
            "org": "TokiStorage",
            "product": "TokiStorage",
        },
        "s2": {
            "bar": "AIは「最適解」を出す。しかし「何を残すか」は決められない。意味の選択は、あなたにしかできない。",
            "label": "Why This Question",
            "cards": [
                ("01", "AIは答えを出す。しかし「問い」は立てられない。",
                 "生成AIは調査・分析・戦略を瞬時に出力します。しかし「100年後に何を残すか」——この問いへの答えは、あなたの人生の文脈からしか生まれません。"),
                ("02", "富の継承は解決済み。意味の継承は手つかず。",
                 "相続設計、事業承継、資産運用——「何を渡すか」の仕組みは整っています。しかし「なぜそれを残すのか」という物語は、制度では設計できません。"),
                ("03", "存在証明は、未来への贈り物である。",
                 "あなたが残すものは、未来の誰かへの贈り物になります。家族に、地域に、まだ生まれていない世代に。「自分が確かに存在した」という証が、未来の誰かの支えになる。"),
            ],
            "footer": "Timeless Consulting",
        },
        "s3": {
            "bar": "あなたが買うのは「プロダクト」ではない。千年の視座を持つ人間との対話である。",
            "label": "Our Approach",
            "quote": "千年の視座を持つ人間との対話から、\nあなたの存在証明を設計します。",
            "cards": [
                ("01", "千年の問いとの対話",
                 "「100年後に何を残すか」——心理学・宗教・経済・AI・宇宙まで9領域の思想フレームワークで向き合います。"),
                ("02", "存在証明のデザイン",
                 "あなたの物語・価値観・メッセージを構造化し、千年先まで届く「存在証明」として設計します。"),
                ("03", "石英ガラスへの刻印",
                 "サーバー・電源・制度に依存しない、1000年保証の記録媒体。スマホカメラだけで読み取り可能。"),
                ("04", "継続的なキュレーション",
                 "人生は一度で語りきれません。節目ごとに存在証明をアップデートし、千年アーカイブを育てます。"),
            ],
            "footer": "Timeless Consulting",
        },
        "s4": {
            "bar": "対話から始まり、千年に届く。4つのステップであなたの存在証明が形になる。",
            "label": "Process",
            "steps": [
                ("01", "初回対話（90分）",
                 "「100年後に何を残したいか」を軸に、あなたの人生の物語を聴かせていただきます。この対話自体が、深い内省の機会になります。"),
                ("02", "存在証明の設計",
                 "対話をもとに、何を・誰に・どんな形で残すかを設計します。テキスト、写真、音声、映像——最適なメディアを提案。"),
                ("03", "コンテンツの制作・キュレーション",
                 "設計に基づいて存在証明を制作。必要に応じてプロフェッショナルな撮影・収録・編集をコーディネート。"),
                ("04", "石英ガラスへの刻印・納品",
                 "完成した存在証明を石英ガラスに刻印し、QRコードとともにお届け。1000年先の誰かがあなたの物語に出会えます。"),
            ],
            "footer": "Timeless Consulting",
        },
        "s5": {
            "bar": "「残す」ことに意味を感じるすべての方へ——6つの領域で特に高い親和性がある",
            "label": "Who This Is For",
            "sectors": [
                ("C", "経営者・創業者", "会社を超えて残る、あなた自身の理念と物語を"),
                ("F", "家族の記録を残す方", "家族の歴史、想い出、伝えたい言葉を千年先へ"),
                ("R", "宗教者・文化人", "教え、作品、精神的遺産を制度に依存せず永続化"),
                ("A", "アーティスト", "作品と創作の意図を、プラットフォームを超えて残す"),
                ("L", "地域・コミュニティ", "まちの記憶、お祭り、方言、災害の教訓を未来に"),
                ("E", "教育者・研究者", "知の系譜、発見の文脈、師弟の物語を永続的に記録"),
            ],
            "footer": "Timeless Consulting",
        },
        "s6": {
            "bar": "対話から生まれる、3つの具体的な成果物。千年先まで届く「かたち」をお届けします。",
            "label": "Deliverables",
            "items": [
                (TOKI_BLUE_PALE, "Q", "石英ガラスの存在証明",
                 "あなたの物語が刻まれた石英ガラス。サーバー・電源不要、1000年保証。スマートフォンのカメラで読み取り可能なQRコード付き。"),
                (GOLD_PALE, "D", "デジタルコンパニオン",
                 "QRコードを読み取ると表示されるデジタルページ。テキスト、写真、音声、映像を含む、存在証明のリッチな体験版。"),
                (GREEN_PALE, "R", "思想フレームワークレポート",
                 "あなたの存在証明を、70以上の思想エッセイの文脈に位置づけたレポート。「なぜ残すのか」の知的な裏付け。"),
            ],
            "footer": "Timeless Consulting",
        },
        "s7": {
            "bar": "千年の視座であなたの物語に向き合う——対話のパートナーについて",
            "label": "Your Partner",
            "name": "佐藤卓也 \u2014 TokiStorage 代表",
            "bio": "大手コンサルティングファームでの経験を経て、半導体製造装置のエンジニアリング20年超。タイムレスタウン新浦安（250世帯）の自治会長として「ゆりかごから墓場まで」のコミュニティ運営を経験。SoulCarrier活動で「記憶が消える恐怖」を目の当たりにし、TokiStorageを着想。マウイ・山中湖でのオフグリッド実証を経て、制度に依存しない千年設計の技術を完成。70以上の思想エッセイを執筆し、9つの知的領域から存在証明の意味を探究し続けている。",
            "tags": ["元Big4ファーム", "半導体エンジニアリング 20年+", "自治会長（250世帯）",
                     "SoulCarrier主宰", "70+思想エッセイ執筆", "オフグリッド実証済み", "佐渡島移住予定（2026春）"],
            "story": "愛犬パールのお墓とともに、家族でパールハーバーを訪れました。平和を願う声を音声QRとして刻んだプレートを持って。そこで気がつきました——時を超えた願いは、愛や平和への想いに変容する。一人でも多くの方と分かち合いたい。TokiStorageをきっかけに、時を超えた変容をご一緒できれば幸いです。",
            "footer": "Timeless Consulting",
        },
        "s8": {
            "bar": "Next Step",
            "title": "まずは、お話しましょう。",
            "steps": [
                ("01", "初回対話（90分・無料）", "「100年後に何を残したいか」を一緒に考える、最初の90分"),
                ("02", "存在証明の設計ご提案", "対話をもとに、あなただけのプランをお作りします"),
                ("03", "制作・刻印", "存在証明を制作し、石英ガラスに永久に刻みます"),
                ("04", "納品・キュレーション開始", "お届けし、継続的なアップデートをサポートします"),
            ],
            "contact": "TokiStorage　佐藤卓也",
            "footer_left": "TokiStorage",
        },
        "s9": {
            "title": "Confidential / Disclaimer",
            "body": "本資料は、TokiStorage（佐藤卓也）がサービスご検討のために作成した資料です。\n\n本資料に含まれる情報は、現時点における見解および計画に基づくものであり、その正確性、完全性、または将来の結果を保証するものではありません。\n\n本資料は情報提供を目的としており、法的助言、投資助言、その他いかなる専門的助言を構成するものでもありません。",
            "copyright": "\u00a9 2026 TokiStorage / 佐藤卓也. All rights reserved.",
            "footer_left": "TokiStorage",
        },
    },
    "en": {
        "font": FONT_EN,
        "filename": "tokistorage-client-deck-en",
        "cover": {
            "label": "Timeless Consulting",
            "title": "What will you preserve\nfor 100 years?",
            "sub": "In an age where AI answers every question,\nthis is the one only you can answer.",
            "org": "TokiStorage",
            "product": "TokiStorage",
        },
        "s2": {
            "bar": "AI generates optimal answers. But it cannot decide what to preserve. The choice of meaning belongs to you alone.",
            "label": "Why This Question",
            "cards": [
                ("01", "AI produces answers. It cannot frame the question.",
                 "Generative AI produces research, analysis, and strategy in seconds. Yet \"What will you preserve for 100 years?\" can only be answered from the context of your own life."),
                ("02", "Wealth transfer is solved. Meaning transfer is not.",
                 "Estate planning, business succession, asset management\u2014the infrastructure for transferring what you own is mature. But the story of why you leave it cannot be designed by institutions."),
                ("03", "Proof of existence is a gift to the future.",
                 "What you leave behind becomes a gift to someone in the future. To your family, your community, generations not yet born. The knowledge that someone was truly here becomes an anchor."),
            ],
            "footer": "Timeless Consulting",
        },
        "s3": {
            "bar": "What you're buying is not a product. It's dialogue with a person who holds a millennium perspective.",
            "label": "Our Approach",
            "quote": "From dialogue with a millennium perspective,\nwe design your proof of existence.",
            "cards": [
                ("01", "Dialogue with the millennium question",
                 "\"What will you preserve for 100 years?\" \u2014 approached through 9 intellectual domains: psychology, religion, economics, AI, and space."),
                ("02", "Proof of existence design",
                 "Your life story, values, and messages structured into a \"proof of existence\" designed to reach a millennium into the future."),
                ("03", "Quartz glass inscription",
                 "No servers, no power, no institutional dependency. Guaranteed for 1,000 years. Readable with any smartphone camera."),
                ("04", "Ongoing curation",
                 "A life cannot be told in a single sitting. At each milestone, we update your proof of existence and grow your archive."),
            ],
            "footer": "Timeless Consulting",
        },
        "s4": {
            "bar": "From dialogue to a millennium. Four steps to give your proof of existence tangible form.",
            "label": "Process",
            "steps": [
                ("01", "Initial Dialogue (90 min)",
                 "Centered on \"What do you want to preserve for 100 years?\", we listen to the story of your life. For many, this dialogue becomes a profound moment of reflection."),
                ("02", "Proof of Existence Design",
                 "Based on our dialogue, we design what to preserve, for whom, in what form. Text, photos, audio, video \u2014 the optimal media and structure."),
                ("03", "Content Creation & Curation",
                 "We produce your proof of existence. Where needed, we coordinate professional photography, recording, and editing."),
                ("04", "Quartz Glass Inscription & Delivery",
                 "Inscribed on quartz glass with a QR code. A smartphone camera is all it takes for someone 1,000 years from now to encounter your story."),
            ],
            "footer": "Timeless Consulting",
        },
        "s5": {
            "bar": "For anyone who finds meaning in leaving something behind \u2014 six high-affinity audiences",
            "label": "Who This Is For",
            "sectors": [
                ("C", "Founders & CEOs", "Your philosophy and story, beyond the company you built"),
                ("F", "Family Legacy Seekers", "Family history, memories, words to pass on for a millennium"),
                ("R", "Religious & Cultural Leaders", "Teachings, works, spiritual heritage beyond institutions"),
                ("A", "Artists & Creators", "Your work and creative intent, beyond any platform"),
                ("L", "Communities", "Town memories, festivals, dialects, disaster lessons for the future"),
                ("E", "Educators & Researchers", "Intellectual lineage, discovery context, mentor-student narratives"),
            ],
            "footer": "Timeless Consulting",
        },
        "s6": {
            "bar": "Three tangible outcomes from a meaningful dialogue. Forms that reach a millennium.",
            "label": "Deliverables",
            "items": [
                (TOKI_BLUE_PALE, "Q", "Quartz Glass Proof of Existence",
                 "Your story inscribed in quartz glass. No servers, no power needed. Guaranteed for 1,000 years. QR code readable by any smartphone."),
                (GOLD_PALE, "D", "Digital Companion",
                 "A digital page accessed via QR code scan. A rich experience including text, photos, audio, and video of your proof of existence."),
                (GREEN_PALE, "R", "Philosophical Framework Report",
                 "A report placing your proof of existence in the context of 70+ philosophical essays. Intellectual grounding for \"why preserve.\""),
            ],
            "footer": "Timeless Consulting",
        },
        "s7": {
            "bar": "The person who engages your story with a millennium perspective \u2014 your dialogue partner",
            "label": "Your Partner",
            "name": "Takuya Sato \u2014 Founder, TokiStorage",
            "bio": "Former Big Four consultant \u2014 understands executive dialogue and strategic thinking from the inside. 20+ years in semiconductor manufacturing engineering. Former president of Timeless Town Shin-Urayasu residents' association (250 households). Through SoulCarrier's work with unclaimed graves, witnessed firsthand how memories vanish \u2014 and conceived TokiStorage. Validated off-grid, institution-free 1,000-year design through testing in Maui and Lake Yamanakako. Author of 70+ philosophical essays exploring proof of existence across 9 intellectual domains.",
            "tags": ["Big Four Alumni", "Semiconductor engineering 20+ yrs", "Community president (250 households)",
                     "SoulCarrier founder", "70+ philosophical essays", "Off-grid validated", "Relocating to Sado Island (Spring 2026)"],
            "story": "I visited Pearl Harbor with my family, carrying a plate inscribed with a voice QR code of our prayers for peace\u2014alongside the remains of our beloved dog, Pearl. That's when it hit me: wishes that transcend time transform into something greater\u2014love, peace, hope. I want to share that transformation with as many people as possible. Through TokiStorage, I hope to walk that journey across time together with you.",
            "footer": "Timeless Consulting",
        },
        "s8": {
            "bar": "Next Step",
            "title": "Let's start with a conversation.",
            "steps": [
                ("01", "Initial Dialogue (90 min, complimentary)", "Explore \"What will you preserve for 100 years?\" together"),
                ("02", "Design Proposal", "A personalized proof of existence plan based on our dialogue"),
                ("03", "Creation & Inscription", "Produce and permanently inscribe in quartz glass"),
                ("04", "Delivery & Curation", "Receive your artifact and begin ongoing curation support"),
            ],
            "contact": "Takuya Sato \u2014 Founder, TokiStorage",
            "footer_left": "TokiStorage",
        },
        "s9": {
            "title": "Confidential / Disclaimer",
            "body": "This document has been prepared by TokiStorage (Takuya Sato) for the purpose of introducing our consulting services.\n\nThe information contained herein reflects current views and plans and does not constitute a guarantee of accuracy, completeness, or future outcomes.\n\nThis document is provided for informational purposes only and does not constitute legal, investment, or other professional advice.",
            "copyright": "\u00a9 2026 TokiStorage / Takuya Sato. All rights reserved.",
            "footer_left": "TokiStorage",
        },
    },
}


# ══════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════

def build_cover(prs, d):
    slide = add_blank_slide(prs)
    font = d["font"]
    c = d["cover"]
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=DARK_BG)
    # Subtle gradient effect via darker bottom strip
    add_rect(slide, 0, SLIDE_H - Inches(1.5), SLIDE_W, Inches(1.5), fill=DARK_BG2)
    add_textbox(slide, Inches(1), Inches(0.7), Inches(5), Inches(0.35),
                c["label"], font, 11, GOLD, align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(1), Inches(1.4), Inches(8), Inches(1.5),
                c["title"], font, 32, WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(3.2), Inches(7), Inches(0.9),
                c["sub"], font, 14, RGBColor(0xBB, 0xBB, 0xCC))
    # Bottom accent
    stripe_y = SLIDE_H - Inches(0.65)
    add_rect(slide, 0, stripe_y, SLIDE_W, Inches(0.03), fill=GOLD)
    add_textbox(slide, Inches(0.5), stripe_y + Inches(0.1), Inches(3), Inches(0.35),
                c["org"], font, 9, TEXT_MUTED)
    add_textbox(slide, Inches(4), stripe_y + Inches(0.1), Inches(2), Inches(0.35),
                c["product"], font, 9, TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, SLIDE_W - Inches(2), stripe_y + Inches(0.1), Inches(1.5), Inches(0.35),
                "Confidential", font, 9, TEXT_MUTED, align=PP_ALIGN.RIGHT)


def build_slide2(prs, d):
    """Why This Question: 3 column cards"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s2"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))
    card_w = Inches(2.85)
    card_h = Inches(3.7)
    gap = Inches(0.23)
    start_x = Inches(0.5)
    y = Inches(1.15)
    for i, (num, title, body) in enumerate(s["cards"]):
        x = start_x + i * (card_w + gap)
        draw_col_card(slide, x, y, card_w, card_h, num, title, body, font)
    add_footer(slide, s["footer"], 2, font)


def build_slide3(prs, d):
    """Our Approach: quote + 2x2 grid"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s3"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))
    # Quote callout
    qx, qy = Inches(0.5), Inches(1.1)
    qw, qh = Inches(8.6), Inches(0.75)
    add_rect(slide, qx, qy, qw, qh, fill=TOKI_BLUE_PALE)
    add_rect(slide, qx, qy, Inches(0.06), qh, fill=TOKI_BLUE)
    add_textbox(slide, qx + Inches(0.25), qy + Inches(0.05), qw - Inches(0.35), qh - Inches(0.1),
                s["quote"], font, 11, TEXT_PRIMARY, bold=True, anchor=MSO_ANCHOR.MIDDLE)
    # 2x2 grid
    card_w = Inches(4.18)
    card_h = Inches(1.4)
    gap_x = Inches(0.25)
    gap_y = Inches(0.12)
    start_x = Inches(0.5)
    start_y = Inches(1.95)
    for i, (num, title, body) in enumerate(s["cards"]):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        draw_grid_card(slide, x, y, card_w, card_h, num, title, body, font)
    add_footer(slide, s["footer"], 3, font)


def build_slide4(prs, d):
    """Process: 4 steps vertical"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s4"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))
    step_w = Inches(8.6)
    step_h = Inches(0.96)
    gap = Inches(0.06)
    start_y = Inches(1.15)
    x = Inches(0.5)
    for i, (num, title, desc) in enumerate(s["steps"]):
        y = start_y + i * (step_h + gap)
        add_rect(slide, x, y, step_w, step_h, fill=WHITE, border_color=BORDER)
        # Number circle
        cx, cy = x + Inches(0.18), y + Inches(0.2)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, cx, cy, Inches(0.5), Inches(0.5))
        circ.fill.solid()
        circ.fill.fore_color.rgb = TOKI_BLUE
        circ.line.fill.background()
        add_textbox(slide, cx, cy, Inches(0.5), Inches(0.5),
                    num, font, 12, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        # Title
        add_textbox(slide, x + Inches(0.85), y + Inches(0.08), step_w - Inches(1.0), Inches(0.3),
                    title, font, 12, TEXT_PRIMARY, bold=True)
        # Description
        add_textbox(slide, x + Inches(0.85), y + Inches(0.36), step_w - Inches(1.0), Inches(0.52),
                    desc, font, 10, TEXT_SECONDARY)
    add_footer(slide, s["footer"], 4, font)


def build_slide5(prs, d):
    """Who This Is For: 3x2 grid"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s5"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))
    card_w = Inches(2.85)
    card_h = Inches(1.55)
    gap_x = Inches(0.23)
    gap_y = Inches(0.15)
    start_x = Inches(0.5)
    start_y = Inches(1.25)
    for i, (icon, title, body) in enumerate(s["sectors"]):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        draw_sector_card(slide, x, y, card_w, card_h, icon, title, body, font)
    add_footer(slide, s["footer"], 5, font)


def build_slide6(prs, d):
    """Deliverables: 3 horizontal items"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s6"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))
    item_w = Inches(8.6)
    item_h = Inches(1.1)
    gap = Inches(0.15)
    start_y = Inches(1.2)
    x = Inches(0.5)
    for i, (bg, icon, title, body) in enumerate(s["items"]):
        y = start_y + i * (item_h + gap)
        draw_deliverable_item(slide, x, y, item_w, item_h, bg, icon, title, body, font)
    add_footer(slide, s["footer"], 6, font)


def build_slide7(prs, d):
    """Founder profile"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s7"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))

    # Compact mode for EN with story (bio is longer, needs space savings)
    has_story = "story" in s
    compact = (font == FONT_EN and has_story)

    # Photo (circular, 1.0")
    ax, ay, asize = Inches(0.5), Inches(1.12), Inches(1.0)
    photo_path = os.path.join(OUT_DIR, "asset", "IMG_4310-2.jpeg")
    if os.path.exists(photo_path):
        try:
            circ_path = _prepare_circular_photo(photo_path)
            slide.shapes.add_picture(circ_path, ax, ay, asize, asize)
        except Exception:
            _draw_avatar_fallback(slide, ax, ay, asize, d, font)
    else:
        _draw_avatar_fallback(slide, ax, ay, asize, d, font)

    # Name
    name_x = Inches(1.65)
    add_textbox(slide, name_x, Inches(1.18), Inches(7.3), Inches(0.35),
                s["name"], font, 13, TEXT_PRIMARY, bold=True)
    # Bio — 9pt in compact mode to save vertical space
    bio_y = 1.58
    bio_size = 9 if compact else 10
    add_textbox(slide, name_x, Inches(bio_y), Inches(7.3), Inches(1.2),
                s["bio"], font, bio_size, TEXT_SECONDARY)

    # Tags (dynamic positioning)
    cpl = 48 if font == FONT_JP else (95 if compact else 85)
    bio_lines = (len(s["bio"]) + cpl - 1) // cpl
    bio_h_est = bio_lines * (0.17 if compact else 0.19)
    tag_start_x = Inches(0.5)
    tag_x = tag_start_x
    tag_y = Inches(max(2.6, bio_y + bio_h_est + 0.30))
    tag_h = Inches(0.38 if compact else 0.44)
    tag_row_gap = Inches(0.42 if compact else 0.52)
    for tag in s["tags"]:
        tw = Inches(len(tag) * 0.075 + 0.45)
        if tag_x + tw > Inches(9.3):
            tag_x = tag_start_x
            tag_y += tag_row_gap
        add_rect(slide, tag_x, tag_y, tw, tag_h, fill=BG_SECTION, border_color=BORDER)
        add_textbox(slide, tag_x + Inches(0.1), tag_y, tw - Inches(0.2), tag_h,
                    tag, font, 8, TEXT_SECONDARY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tag_x += tw + Inches(0.12)

    # Story callout (emotional origin) — dynamically sized to fit
    if has_story:
        footer_y = SLIDE_H - Inches(0.38)
        sy = tag_y + Inches(0.56)
        available = footer_y - sy - Inches(0.1)
        sh = min(Inches(0.85), available)
        if sh >= Inches(0.45):
            story_font = 9 if sh >= Inches(0.65) else 8
            sx, sw = Inches(0.5), Inches(8.6)
            add_rect(slide, sx, sy, sw, sh, fill=TOKI_BLUE_PALE)
            add_rect(slide, sx, sy, Inches(0.06), sh, fill=GOLD)
            add_textbox(slide, sx + Inches(0.25), sy + Inches(0.08), sw - Inches(0.35), sh - Inches(0.16),
                        s["story"], font, story_font, TEXT_PRIMARY, bold=False)

    add_footer(slide, s["footer"], 7, font)


def _draw_avatar_fallback(slide, ax, ay, asize, d, font):
    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, ax, ay, asize, asize)
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = DARK_BG
    avatar.line.fill.background()
    initials = "佐" if d is CONTENT["ja"] else "TS"
    add_textbox(slide, ax, ay, asize, asize,
                initials, font, 16, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def build_slide8(prs, d):
    """Next Steps — dark background"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s8"]
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=DARK_BG)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(4), Inches(0.45),
                s["bar"], font, 13, WHITE, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(8), Inches(0.5),
                s["title"], font, 20, WHITE, bold=True)
    # 4 numbered steps
    step_y = Inches(1.6)
    step_h = Inches(0.72)
    step_gap = Inches(0.06)
    for i, (num, title, desc) in enumerate(s["steps"]):
        y = step_y + i * (step_h + step_gap)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(0.8), y + Inches(0.12),
                                       Inches(0.42), Inches(0.42))
        circ.fill.solid()
        circ.fill.fore_color.rgb = TOKI_BLUE
        circ.line.fill.background()
        add_textbox(slide, Inches(0.8), y + Inches(0.12), Inches(0.42), Inches(0.42),
                    num, font, 10, WHITE, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.4), y + Inches(0.06), Inches(7), Inches(0.3),
                    title, font, 12, WHITE, bold=True)
        add_textbox(slide, Inches(1.4), y + Inches(0.38), Inches(7), Inches(0.3),
                    desc, font, 10, RGBColor(0xBB, 0xBB, 0xCC))
    # Contact
    add_textbox(slide, Inches(0.8), SLIDE_H - Inches(1.0), Inches(8), Inches(0.3),
                s["contact"], font, 10, TEXT_MUTED)
    # Footer
    stripe_y = SLIDE_H - Inches(0.55)
    add_rect(slide, 0, stripe_y - Inches(0.03), SLIDE_W, Pt(0.5), fill=RGBColor(0x33, 0x44, 0x55))
    add_textbox(slide, Inches(0.5), stripe_y, Inches(4), Inches(0.35),
                s["footer_left"], font, 9, TEXT_MUTED)
    add_textbox(slide, Inches(4), stripe_y, Inches(2), Inches(0.35),
                "Confidential", font, 9, TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, SLIDE_W - Inches(1), stripe_y, Inches(0.5), Inches(0.35),
                "8", font, 9, TEXT_MUTED, bold=True, align=PP_ALIGN.RIGHT)


def build_slide9(prs, d):
    """Disclaimer"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s9"]
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG_PAGE)
    add_textbox(slide, Inches(1), Inches(1.0), Inches(8), Inches(0.45),
                s["title"], font, 14, TEXT_MUTED, bold=True)
    add_rect(slide, Inches(1), Inches(1.5), Inches(8), Pt(0.5), fill=BORDER)
    add_textbox(slide, Inches(1), Inches(1.7), Inches(8), Inches(2.8),
                s["body"], font, 9, TEXT_SECONDARY)
    add_textbox(slide, Inches(1), SLIDE_H - Inches(1.2), Inches(8), Inches(0.3),
                s["copyright"], font, 8, TEXT_MUTED)
    add_footer(slide, s["footer_left"], 9, font)


# ══════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════

def generate(lang):
    d = CONTENT[lang]
    prs = new_prs()
    build_cover(prs, d)
    build_slide2(prs, d)
    build_slide3(prs, d)
    build_slide4(prs, d)
    build_slide5(prs, d)
    build_slide6(prs, d)
    build_slide7(prs, d)
    build_slide8(prs, d)
    build_slide9(prs, d)

    pptx_path = os.path.join(OUT_DIR, f"{d['filename']}.pptx")
    prs.save(pptx_path)
    _strip_theme_shadows(pptx_path)
    print(f"  PPTX saved: {pptx_path}")
    return pptx_path


def _strip_theme_shadows(pptx_path):
    clean = '<a:effectStyleLst>' + \
            '<a:effectStyle><a:effectLst/></a:effectStyle>' * 3 + \
            '</a:effectStyleLst>'
    buf = io.BytesIO()
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                data = zin.read(item)
                if item == 'ppt/theme/theme1.xml':
                    text = data.decode('utf-8')
                    text = re.sub(r'<a:effectStyleLst>.*?</a:effectStyleLst>',
                                  clean, text, flags=re.DOTALL)
                    data = text.encode('utf-8')
                zout.writestr(item, data)
    with open(pptx_path, 'wb') as f:
        f.write(buf.getvalue())


def convert_to_pdf(pptx_path):
    out_dir = os.path.dirname(pptx_path)
    env = os.environ.copy()
    env["HOME"] = "/tmp"
    result = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
        capture_output=True, text=True, timeout=120, env=env
    )
    if result.returncode != 0:
        print(f"  ERROR: {result.stderr}")
        return None
    pdf_path = pptx_path.replace(".pptx", ".pdf")
    print(f"  PDF saved: {pdf_path}")
    return pdf_path


if __name__ == "__main__":
    print("=== TokiStorage Client Deck Generator (Timeless Consulting) ===\n")

    for lang in ["ja", "en"]:
        print(f"[{lang.upper()}] Generating PPTX...")
        pptx = generate(lang)
        print(f"[{lang.upper()}] Converting to PDF...")
        convert_to_pdf(pptx)

    print("\nDone!")
