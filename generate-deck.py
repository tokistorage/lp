#!/usr/bin/env python3
"""
TokiStorage Partnership Deck Generator
Generates PPTX files (JP + EN) and converts to PDF via LibreOffice.
Design aligned with TokiStorage landing page (index.css).
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import subprocess, os, sys, re, zipfile, io, tempfile

# ── Circular photo helper ─────────────────────────────────────────────

def _prepare_circular_photo(src_path, size=300):
    """Crop photo to circle with transparent background. Returns temp PNG path."""
    from PIL import Image, ImageDraw
    img = Image.open(src_path).convert("RGBA")
    w, h = img.size
    s = min(w, h)
    left, top = (w - s) // 2, (h - s) // 2
    img = img.crop((left, top, left + s, top + s)).resize((size, size), Image.LANCZOS)
    mask = Image.new("L", (size, size), 0)
    ImageDraw.Draw(mask).ellipse((0, 0, size, size), fill=255)
    result = Image.new("RGBA", (size, size), (0, 0, 0, 0))
    result.paste(img, mask=mask)
    tmp = os.path.join(tempfile.gettempdir(), "_deck_profile.png")
    result.save(tmp, "PNG")
    return tmp

# ── Design tokens (matched to index.css :root) ────────────────────────
TOKI_BLUE     = RGBColor(0x25, 0x63, 0xEB)  # --toki-blue
TOKI_BLUE_DK  = RGBColor(0x1D, 0x4E, 0xD8)  # --toki-blue-dark
TOKI_BLUE_PALE= RGBColor(0xEF, 0xF6, 0xFF)  # --toki-blue-pale
GOLD          = RGBColor(0xC9, 0xA9, 0x62)  # --toki-gold
GOLD_PALE     = RGBColor(0xFD, 0xF8, 0xE8)  # --toki-gold-pale
EMERALD       = RGBColor(0x16, 0xA3, 0x4A)  # green-600 (check marks in LP)
GREEN_PALE    = RGBColor(0xF0, 0xFD, 0xF4)  # green-50

TEXT_PRIMARY  = RGBColor(0x1E, 0x29, 0x3B)  # --text-primary
TEXT_SECONDARY= RGBColor(0x47, 0x55, 0x69)  # --text-secondary
TEXT_MUTED    = RGBColor(0x94, 0xA3, 0xB8)  # --text-muted
WHITE         = RGBColor(0xFF, 0xFF, 0xFF)

BG_PAGE       = RGBColor(0xF8, 0xFA, 0xFC)  # --bg-page
BG_SECTION    = RGBColor(0xF1, 0xF5, 0xF9)  # --bg-section
BORDER        = RGBColor(0xE2, 0xE8, 0xF0)  # --border

# Dark backgrounds (LP hero / footer)
DARK_BG       = RGBColor(0x1E, 0x29, 0x3B)  # --text-primary as bg (LP footer)
DARK_BG2      = RGBColor(0x0F, 0x17, 0x2A)  # LP footer gradient start

SLIDE_W = Inches(10)
SLIDE_H = Inches(5.625)  # 16:9

FONT_JP = "IPAPGothic"
FONT_EN = "Calibri"

OUT_DIR = os.path.dirname(os.path.abspath(__file__))


# ── Helpers ────────────────────────────────────────────────────────────

def new_prs():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def add_blank_slide(prs):
    layout = prs.slide_layouts[6]  # blank
    return prs.slides.add_slide(layout)


def add_rect(slide, left, top, width, height, fill=None, border_color=None, border_width=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    if fill is not None:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    if border_color:
        shape.line.color.rgb = border_color
        shape.line.width = border_width or Pt(0.75)
    else:
        shape.line.fill.background()
    return shape


def set_text(tf, text, font_name, size, color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, line_spacing=1.35):
    tf.word_wrap = True
    for p in tf.paragraphs:
        p.clear()
    p = tf.paragraphs[0]
    p.alignment = align
    p.line_spacing = line_spacing
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_para(tf, text, font_name, size, color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, space_before=0, line_spacing=1.35):
    p = tf.add_paragraph()
    p.alignment = align
    p.line_spacing = line_spacing
    if space_before:
        p.space_before = Pt(space_before)
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(size)
    run.font.color.rgb = color
    run.font.bold = bold
    return p


def add_textbox(slide, left, top, width, height, text, font_name, size,
                color=TEXT_PRIMARY, bold=False, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.auto_size = None
    try:
        txBox.text_frame.paragraphs[0].space_before = Pt(0)
        txBox.text_frame.paragraphs[0].space_after = Pt(0)
    except:
        pass
    from pptx.oxml.ns import qn
    txBody = txBox.text_frame._txBody
    bodyPr = txBody.find(qn('a:bodyPr'))
    if anchor == MSO_ANCHOR.MIDDLE:
        bodyPr.set('anchor', 'ctr')
    elif anchor == MSO_ANCHOR.BOTTOM:
        bodyPr.set('anchor', 'b')
    set_text(tf, text, font_name, size, color, bold, align)
    return txBox


# ── Action bar ─────────────────────────────────────────────────────────

def add_action_bar(slide, text, font):
    bar_h = Inches(0.65)
    add_rect(slide, 0, 0, SLIDE_W, bar_h, fill=DARK_BG)
    add_textbox(slide, Inches(0.5), 0, SLIDE_W - Inches(1), bar_h,
                text, font, 12, WHITE, bold=True, anchor=MSO_ANCHOR.MIDDLE)


# ── Footer ─────────────────────────────────────────────────────────────

def add_footer(slide, left_text, pg, font):
    y = SLIDE_H - Inches(0.38)
    add_rect(slide, 0, y, SLIDE_W, Pt(0.5), fill=BORDER)
    add_textbox(slide, Inches(0.5), y + Pt(2), Inches(4), Inches(0.3),
                left_text, font, 9, TEXT_MUTED)
    add_textbox(slide, Inches(4), y + Pt(2), Inches(2), Inches(0.3),
                "Confidential", font, 9, TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, SLIDE_W - Inches(1), y + Pt(2), Inches(0.5), Inches(0.3),
                str(pg), font, 9, TEXT_MUTED, bold=True, align=PP_ALIGN.RIGHT)


# ── Section label ──────────────────────────────────────────────────────

def add_section_label(slide, text, font, top):
    add_textbox(slide, Inches(0.5), top, Inches(3), Inches(0.3),
                text.upper(), font, 10, TOKI_BLUE, bold=True)


# ── Card helpers (all use single border style: BORDER, 0.75pt) ────────

def draw_col_card(slide, x, y, w, h, num, title, body, font):
    add_rect(slide, x, y, w, h, fill=WHITE, border_color=BORDER)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.1), Inches(0.6), Inches(0.35),
                num, font, 15, TOKI_BLUE, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(0.45), w - Inches(0.3), Inches(0.5),
                title, font, 12, TEXT_PRIMARY, bold=True)
    add_textbox(slide, x + Inches(0.15), y + Inches(1.05), w - Inches(0.3), h - Inches(1.15),
                body, font, 12, TEXT_SECONDARY)


def draw_grid_card(slide, x, y, w, h, icon_letter, title, body, font):
    add_rect(slide, x, y, w, h, fill=WHITE, border_color=BORDER)
    # Icon circle
    ix, iy = x + Inches(0.15), y + Inches(0.15)
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, ix, iy, Inches(0.4), Inches(0.4))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TOKI_BLUE_PALE
    circ.line.fill.background()
    add_textbox(slide, ix, iy, Inches(0.4), Inches(0.4),
                icon_letter, font, 12, TOKI_BLUE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, x + Inches(0.65), y + Inches(0.15), w - Inches(0.8), Inches(0.45),
                title, font, 12, TEXT_PRIMARY, bold=True)
    add_textbox(slide, x + Inches(0.65), y + Inches(0.72), w - Inches(0.8), h - Inches(0.82),
                body, font, 11, TEXT_SECONDARY)


def draw_model_item(slide, x, y, w, h, badge_text, badge_color, title, body, example, font):
    add_rect(slide, x, y, w, h, fill=WHITE, border_color=BORDER)
    # Badge (vertically centered in card)
    bx, by, bw, bh = x + Inches(0.15), y + Inches(0.13), Inches(0.85), Inches(0.9)
    badge = add_rect(slide, bx, by, bw, bh, fill=badge_color)
    add_textbox(slide, bx, by, bw, bh, badge_text, font, 8, WHITE, bold=True,
                align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    tx = x + Inches(1.15)
    tw = w - Inches(1.3)
    add_textbox(slide, tx, y + Inches(0.08), tw, Inches(0.38),
                title, font, 11, TEXT_PRIMARY, bold=True)
    add_textbox(slide, tx, y + Inches(0.40), tw, Inches(0.34),
                body, font, 10, TEXT_SECONDARY)
    add_textbox(slide, tx, y + Inches(0.8), tw, Inches(0.22),
                example, font, 9, TOKI_BLUE)


def draw_flow_box(slide, x, y, w, h, title, body, bg_color, title_color, font):
    add_rect(slide, x, y, w, h, fill=bg_color, border_color=BORDER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.08), w - Inches(0.2), Inches(0.3),
                title, font, 11, title_color, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.1), y + Inches(0.46), w - Inches(0.2), h - Inches(0.54),
                body, font, 10, TEXT_SECONDARY, align=PP_ALIGN.CENTER)


def draw_sector_card(slide, x, y, w, h, icon_letter, title, body, font):
    add_rect(slide, x, y, w, h, fill=WHITE, border_color=BORDER)
    # Icon circle
    circ_x = x + (w - Inches(0.42)) / 2
    circ = slide.shapes.add_shape(MSO_SHAPE.OVAL, circ_x, y + Inches(0.1), Inches(0.42), Inches(0.42))
    circ.fill.solid()
    circ.fill.fore_color.rgb = TOKI_BLUE_PALE
    circ.line.fill.background()
    add_textbox(slide, circ_x, y + Inches(0.1), Inches(0.42), Inches(0.42),
                icon_letter, font, 12, TOKI_BLUE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, x + Inches(0.08), y + Inches(0.56), w - Inches(0.16), Inches(0.26),
                title, font, 11, TEXT_PRIMARY, bold=True, align=PP_ALIGN.CENTER)
    add_textbox(slide, x + Inches(0.08), y + Inches(0.9), w - Inches(0.16), h - Inches(0.98),
                body, font, 9, TEXT_SECONDARY, align=PP_ALIGN.CENTER)


def draw_table_row(slide, y, row_h, cells, font, is_header=False, is_first_col_header=False):
    col_widths = [Inches(1.5), Inches(3.5), Inches(3.6)]
    x = Inches(0.7)
    for i, (text, width) in enumerate(zip(cells, col_widths)):
        bg = DARK_BG if is_header else (BG_SECTION if (i == 0 and is_first_col_header) else WHITE)
        fg = WHITE if is_header else (TEXT_PRIMARY if (i == 0 and is_first_col_header) else TEXT_SECONDARY)
        bld = is_header or (i == 0 and is_first_col_header)
        fs = 9 if is_header else 10
        rect = add_rect(slide, x, y, width, row_h, fill=bg, border_color=BORDER, border_width=Pt(0.5))
        if i == 2 and not is_header:
            rect.fill.solid()
            rect.fill.fore_color.rgb = TOKI_BLUE_PALE
            fg = TOKI_BLUE_DK
        add_textbox(slide, x + Inches(0.1), y, width - Inches(0.2), row_h,
                    text, font, fs, fg, bold=bld, anchor=MSO_ANCHOR.MIDDLE)
        x += width


# ══════════════════════════════════════════════════════════════════════
#  CONTENT DATA
# ══════════════════════════════════════════════════════════════════════

CONTENT = {
    "ja": {
        "font": FONT_JP,
        "filename": "tokistorage-partnership-deck",
        "cover": {
            "label": "Partnership Proposal \u2014 Confidential",
            "title": "貴社クライアントに、\n「千年の存在証明」を。",
            "sub": "テクノロジーは揃っています。ユースケースも200以上。\n足りないのは、届ける仕組みです。",
            "org": "TokiStorage",
            "product": "TokiStorage",
        },
        "s2": {
            "bar": "AIとビッグテックの台頭は、コンサルティングファームに「コモディティ化しない問い」の獲得を迫っている",
            "label": "Background",
            "cards": [
                ("01", "AIがコンサルティングを変える",
                 "生成AIにより、調査・分析・戦略フレームワーク作成は急速にコモディティ化。知識の非対称性に依存した価値提供モデルは構造的転換期にある。"),
                ("02", "差別化の源泉が縮小",
                 "業界分析、ベンチマーク、DX——どのファームも同じツールで同じ提案に収束しがち。「何を提案するか」ではなく「どんな問いを立てるか」が差別化の核になる。"),
                ("03", "「千年の問い」という機会",
                 "「100年後に何を残すか」——AIが生成できない問い。クライアントの本質的ニーズに触れ、構造的にコモディティ化しない提案の起点になる。"),
            ],
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s3": {
            "bar": "TokiStorageは技術・ユースケース・思想基盤・展開網の4つを「設計思想ごと」提供できる",
            "label": "Our Offering",
            "cards": [
                ("Q", "石英ガラス記録技術",
                 "金属蒸着によるQRコード刻印。サーバー・電源ゼロ。SLA 100%、1000年保証。スマホカメラで読取可能。"),
                ("U", "200+ ユースケース（業界別整理済み）",
                 "終活・婚礼・寺社・学校・企業・自治体・NGO・ホテル・航空。提案書にそのまま転用可能な粒度。"),
                ("E", "70+ 思想エッセイ（9領域）",
                 "存在証明を心理学・宗教・経済・AI・宇宙まで展開。知的コンテンツとしてクライアント提案に活用可能。"),
                ("A", "Pearl Soap + アンバサダー網",
                 "贈与経済の実践。全国にワークショップ展開可能な分散型運動体。エンドユーザーとの接点を提供。"),
            ],
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s4": {
            "bar": "TokiStorageは既存デジタルサービスの「競合」ではなく「千年レイヤー」として補完する位置づけである",
            "label": "Positioning",
            "headers": ["", "デジタルサービス（日常の記録）", "TokiStorage（千年の記録）"],
            "rows": [
                ("記録媒体", "クラウド / HDD", "石英ガラス（物理）"),
                ("得意な時間軸", "今〜数十年（日常利用に最適）", "100年〜1000年（永続保存に特化）"),
                ("インフラ", "サーバー・電源（利便性の源泉）", "不要（GitHub分散管理）"),
                ("読み取り", "専用アプリ / ログイン", "スマホのカメラだけ"),
                ("文化的厚み", "機能・利便性が価値の中心", "70+ エッセイ＋贈与経済の実践"),
                ("社会的接点", "プラットフォームとしての貢献", "SoulCarrier（無縁墓・遺骨送還活動）"),
            ],
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s5": {
            "bar": "貴社のビジネスモデルに応じた3つの提携モデルを用意しており、段階的な移行も可能",
            "label": "Partnership Models",
            "models": [
                ("紹介\nモデル", TOKI_BLUE,
                 "A. クライアント紹介型パートナーシップ",
                 "紹介ベースで連携。リファラルフィーをお支払い。クライアント対応・納品はTokiStorage側で完結。",
                 "例：終活コンサル→存在証明を提案 / 葬祭業DX→メモリアルオプション追加"),
                ("共同\n提案", GOLD,
                 "B. 共同ソリューション型",
                 "貴社コンサルにTokiStorageを組み込んだ共同提案。ESG・地方創生・文化保存の「出口」として千年記録を位置づけ。",
                 "例：自治体DX→地域記憶アーカイブ / ホテルCX改革→ゲスト記録のアップグレード"),
                ("事業\n共創", EMERALD,
                 "C. 新規事業共創型",
                 "存在証明を軸に新規事業を共同立ち上げ。技術・思想・ユースケースはTokiStorage、市場アクセス・信用・スケールは貴社。",
                 "例：メモリアルテック新規事業 / 企業向け永続アーカイブサービス"),
            ],
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s6": {
            "bar": "クライアント・パートナー・TokiStorageの三者がWinになる収益設計を採用している",
            "label": "Revenue Flow",
            "flows": [
                ("クライアント", "千年の存在証明\n社会的意義の実感", GREEN_PALE, EMERALD),
                ("パートナー（貴社）", "紹介フィー or 共同提案収益\nクライアントLTV向上", TOKI_BLUE_PALE, TOKI_BLUE),
                ("TokiStorage", "技術・思想・納品\n収益の一部→SoulCarrier活動", GOLD_PALE, RGBColor(0x92, 0x40, 0x0E)),
            ],
            "callout_title": "初期パートナー優遇",
            "callout_body": "複数のコンサルティングファームに順次ご提案を進めています。最初に提携いただいたファームには、紹介モデルの優先条件・エリア独占権など、初期パートナーならではのインセンティブをご用意します。",
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s7": {
            "bar": "葬祭・ホスピタリティ・宗教法人・自治体・ESG・金融の6領域で特に高い親和性がある",
            "label": "Client Fit",
            "lead": "貴社のクライアントポートフォリオに、以下のセクターはありませんか。",
            "sectors": [
                ("M", "葬祭・メモリアル", "墓じまい代替、永代供養デジタル化、遺族向け新サービス"),
                ("H", "ホスピタリティ", "ウェディング記録、ホテルCX、記念日サービス"),
                ("R", "宗教法人・寺社", "檀家記録の永続化、参拝体験DX、文化財保存"),
                ("G", "自治体・教育", "地域アーカイブ、災害記録、学校史の永続化"),
                ("E", "ESG・サステナビリティ", "1000年設計の企業理念記録、SDGs実績の永続証明"),
                ("F", "金融・保険", "終活関連サービス連携、デジタル遺品対策"),
            ],
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s8": {
            "bar": "代表はBig4出身であり、ファームのコンプライアンス要件を理解した上で提携モデルを設計している",
            "label": "Team & Independence",
            "name": "佐藤卓也 \u2014 TokiStorage 代表",
            "bio": "大手コンサルティングファームでの経験を経て、半導体製造装置のエンジニアリング20年超。タイムレスタウン新浦安（250世帯）の自治会長として「ゆりかごから墓場まで」のコミュニティ運営を経験。SoulCarrier活動で「記憶が消える恐怖」を目の当たりにし、TokiStorageを着想。マウイ・山中湖でのオフグリッド実証を経て、制度に依存しない千年設計の技術を完成。",
            "tags": ["元Big4ファーム", "半導体エンジニアリング 20年+", "自治会長（250世帯）",
                     "SoulCarrier主宰", "オフグリッド実証済み", "佐渡島移住予定（2026春）"],
            "ind_title": "独立性（Independence）について",
            "ind_body": "本提携はベンダーパートナーシップです。SalesforceやSAPの導入推奨と同じ構造であり、監査契約・出資関係は一切含みません。独立性に関する懸念が発生しない設計です。",
            "footer": "TokiStorage \u2014 協業提案",
        },
        "s9": {
            "bar": "Next Step",
            "title": "ご検討のステップ",
            "steps": [
                ("01", "初回ミーティング（30分）", "貴社クライアントの課題感・ポートフォリオを共有"),
                ("02", "ユースケース選定", "貴社クライアントとの親和性が高い領域を整理"),
                ("03", "提携モデル設計", "最適なモデルの選択・条件・スコープの整理"),
                ("04", "パイロット実施", "1〜2件のクライアントで実証・フィードバック"),
            ],
            "contact": "TokiStorage　佐藤卓也",
            "footer_left": "TokiStorage",
        },
        "s10": {
            "title": "Confidential / Disclaimer",
            "body": "本資料は、TokiStorage（佐藤卓也）がパートナーシップのご検討のために作成した機密資料です。\n\n本資料に含まれる情報は、現時点における見解および計画に基づくものであり、その正確性、完全性、または将来の結果を保証するものではありません。\n\n本資料は情報提供を目的としており、法的助言、投資助言、その他いかなる専門的助言を構成するものでもありません。\n\n本資料の全部または一部を、事前の書面による同意なく、第三者に開示、複製、または配布することを禁じます。",
            "copyright": "\u00a9 2026 TokiStorage / 佐藤卓也. All rights reserved.",
            "footer_left": "TokiStorage",
        },
    },
    "en": {
        "font": FONT_EN,
        "filename": "tokistorage-partnership-deck-en",
        "cover": {
            "label": "Partnership Proposal \u2014 Confidential",
            "title": "Bring your clients\n1,000-year proof of existence.",
            "sub": "The technology is ready. Over 200 use cases mapped.\nWhat's missing is the delivery network.",
            "org": "TokiStorage",
            "product": "TokiStorage",
        },
        "s2": {
            "bar": "The rise of AI and Big Tech demands consulting firms find questions that resist commoditization",
            "label": "Background",
            "cards": [
                ("01", "AI is transforming consulting",
                 "Generative AI is rapidly commoditizing research, analysis, and strategy frameworks. Value models built on information asymmetry face structural disruption."),
                ("02", "Differentiation is shrinking",
                 "Industry analysis, benchmarks, digital transformation\u2014every firm converges on similar proposals with similar tools. The differentiator is shifting from what to propose to what questions to ask."),
                ("03", "The millennium question",
                 "\"What will you preserve for 1,000 years?\"\u2014a question AI cannot generate. It touches clients' fundamental needs and creates proposals structurally immune to commoditization."),
            ],
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s3": {
            "bar": "TokiStorage delivers technology, use cases, intellectual foundation, and distribution as a unified design philosophy",
            "label": "Our Offering",
            "cards": [
                ("Q", "Quartz glass recording",
                 "QR codes inscribed via metal deposition. Zero servers, zero power. SLA 100%, guaranteed 1,000 years. Readable by any smartphone camera."),
                ("U", "200+ use cases (organized by industry)",
                 "End-of-life, weddings, temples, schools, corporations, municipalities, NGOs, hotels, airlines. Ready for direct proposal integration."),
                ("E", "70+ philosophical essays (9 domains)",
                 "Proof of existence explored across psychology, religion, economics, AI, and space. Standalone intellectual content for client proposals."),
                ("A", "Pearl Soap + Ambassador network",
                 "A gift-economy practice and decentralized workshop network ready to scale nationwide. Direct end-user touchpoint."),
            ],
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s4": {
            "bar": "TokiStorage is not a \"competitor\" to digital services \u2014 it is a complementary millennium layer",
            "label": "Positioning",
            "headers": ["", "Digital services (everyday records)", "TokiStorage (millennium records)"],
            "rows": [
                ("Medium", "Cloud / HDD", "Quartz glass (physical)"),
                ("Best horizon", "Now to decades (optimized for daily use)", "100\u20131,000 years (optimized for permanence)"),
                ("Infrastructure", "Servers & power (source of convenience)", "None required (GitHub distributed)"),
                ("Reading", "App / login required", "Any smartphone camera"),
                ("Cultural depth", "Functionality & convenience at the core", "70+ essays + gift economy practice"),
                ("Social impact", "Platform-level contribution", "SoulCarrier (unclaimed graves mission)"),
            ],
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s5": {
            "bar": "Three partnership models tailored to your business model, with progressive escalation possible",
            "label": "Partnership Models",
            "models": [
                ("Referral", TOKI_BLUE,
                 "A. Client Referral Partnership",
                 "Introduce clients when TokiStorage fits. You receive a referral fee; we handle delivery end-to-end.",
                 "E.g.: End-of-life consulting \u2192 offer proof of existence / Funeral DX \u2192 add memorial option"),
                ("Joint", GOLD,
                 "B. Joint Solution Partnership",
                 "Embed TokiStorage into your consulting engagements. Position millennium records as the \"outcome layer\" of ESG, revitalization, or DX projects.",
                 "E.g.: Municipal DX \u2192 community archive / Hotel CX \u2192 guest record upgrade"),
                ("Co-Create", EMERALD,
                 "C. New Business Co-Creation",
                 "Launch a new venture together. We bring technology, philosophy, and use cases. You bring market access, credibility, and scale.",
                 "E.g.: Memorial-tech startup / Enterprise perpetual archive service"),
            ],
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s6": {
            "bar": "Revenue design ensures all three parties \u2014 client, partner, and TokiStorage \u2014 win",
            "label": "Revenue Flow",
            "flows": [
                ("Client", "1,000-year proof of existence\nTangible social meaning", GREEN_PALE, EMERALD),
                ("Partner (you)", "Referral fee or joint revenue\nClient LTV increase", TOKI_BLUE_PALE, TOKI_BLUE),
                ("TokiStorage", "Technology & delivery\nRevenue \u2192 SoulCarrier mission", GOLD_PALE, RGBColor(0x92, 0x40, 0x0E)),
            ],
            "callout_title": "Early Partner Advantage",
            "callout_body": "We are approaching consulting firms sequentially. The first firm to partner receives preferential terms \u2014 including priority referral conditions and potential regional exclusivity. Early movers shape the partnership.",
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s7": {
            "bar": "Six client sectors show particularly high affinity: funeral, hospitality, religious, government, ESG, and finance",
            "label": "Client Fit",
            "lead": "Does your client portfolio include any of these sectors?",
            "sectors": [
                ("M", "Funeral & Memorial", "Gravestone alternatives, digital perpetual care, bereavement services"),
                ("H", "Hospitality", "Wedding records, hotel CX, anniversary services"),
                ("R", "Religious Institutions", "Perpetual congregation records, visitor DX, cultural preservation"),
                ("G", "Government & Education", "Community archives, disaster records, school history"),
                ("E", "ESG & Sustainability", "1,000-year corporate purpose records, SDG impact proof"),
                ("F", "Finance & Insurance", "End-of-life service integration, digital estate"),
            ],
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s8": {
            "bar": "The founder is a Big Four alumnus who designed the partnership model with full awareness of firm compliance",
            "label": "Team & Independence",
            "name": "Takuya Sato \u2014 Founder, TokiStorage",
            "bio": "Former Big Four consultant \u2014 understands firm culture, client engagement, and project design from the inside. 20+ years in semiconductor manufacturing engineering. Former president of Timeless Town Shin-Urayasu residents' association (250 households). Through SoulCarrier's work with unclaimed graves, witnessed firsthand how memories vanish \u2014 and conceived TokiStorage. Validated off-grid, institution-free 1,000-year design through testing in Maui and Lake Yamanakako.",
            "tags": ["Big Four Alumni", "Semiconductor engineering 20+ yrs", "Community president (250 households)",
                     "SoulCarrier founder", "Off-grid validated", "Relocating to Sado Island (Spring 2026)"],
            "ind_title": "A note on independence",
            "ind_body": "This is a vendor partnership \u2014 structurally identical to recommending Salesforce or SAP. No audit engagement, no equity relationship, no independence concerns. Designed with full awareness of firm compliance requirements.",
            "footer": "TokiStorage \u2014 Partnership Proposal",
        },
        "s9": {
            "bar": "Next Step",
            "title": "Proposed Timeline",
            "steps": [
                ("01", "Initial Meeting (30 min)", "Share your client landscape and current challenges"),
                ("02", "Use Case Selection", "Identify high-affinity sectors from your portfolio"),
                ("03", "Partnership Design", "Select model, define scope and terms"),
                ("04", "Pilot Engagement", "Prove value with 1\u20132 client engagements"),
            ],
            "contact": "Takuya Sato \u2014 Founder, TokiStorage",
            "footer_left": "TokiStorage",
        },
        "s10": {
            "title": "Confidential / Disclaimer",
            "body": "This document has been prepared by TokiStorage (Takuya Sato) solely for the purpose of evaluating a potential partnership.\n\nThe information contained herein reflects current views and plans and does not constitute a guarantee of accuracy, completeness, or future outcomes.\n\nThis document is provided for informational purposes only and does not constitute legal, investment, or other professional advice.\n\nNo part of this document may be disclosed, reproduced, or distributed to any third party without prior written consent.",
            "copyright": "\u00a9 2026 TokiStorage / Takuya Sato. All rights reserved.",
            "footer_left": "TokiStorage",
        },
    },
}


# ══════════════════════════════════════════════════════════════════════
#  SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════

def build_cover(prs, d):
    slide = add_blank_slide(prs)
    font = d["font"]
    c = d["cover"]
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=DARK_BG)
    add_textbox(slide, Inches(1), Inches(0.7), Inches(5), Inches(0.35),
                c["label"], font, 10, TEXT_MUTED, align=PP_ALIGN.LEFT)
    add_textbox(slide, Inches(1), Inches(1.3), Inches(8), Inches(1.5),
                c["title"], font, 30, WHITE, bold=True)
    add_textbox(slide, Inches(1), Inches(3.0), Inches(7), Inches(0.9),
                c["sub"], font, 14, RGBColor(0xBB, 0xBB, 0xCC))
    # Bottom accent line
    stripe_y = SLIDE_H - Inches(0.65)
    add_rect(slide, 0, stripe_y, SLIDE_W, Inches(0.03), fill=TOKI_BLUE)
    add_textbox(slide, Inches(0.5), stripe_y + Inches(0.1), Inches(3), Inches(0.35),
                c["org"], font, 9, TEXT_MUTED)
    add_textbox(slide, Inches(4), stripe_y + Inches(0.1), Inches(2), Inches(0.35),
                c["product"], font, 9, TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, SLIDE_W - Inches(2), stripe_y + Inches(0.1), Inches(1.5), Inches(0.35),
                "Confidential", font, 9, TEXT_MUTED, align=PP_ALIGN.RIGHT)


def build_slide2(prs, d):
    """Problem: 3 column cards"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s2"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))
    card_w = Inches(2.85)
    card_h = Inches(3.7)
    gap = Inches(0.23)
    start_x = Inches(0.5)
    y = Inches(1.15)
    for i, (num, title, body) in enumerate(s["cards"]):
        x = start_x + i * (card_w + gap)
        draw_col_card(slide, x, y, card_w, card_h, num, title, body, font)
    add_footer(slide, s["footer"], 2, font)


def build_slide3(prs, d):
    """Solution: 2x2 grid"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s3"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))
    card_w = Inches(4.3)
    card_h = Inches(1.95)
    gap_x = Inches(0.25)
    gap_y = Inches(0.15)
    start_x = Inches(0.5)
    start_y = Inches(1.15)
    for i, (icon, title, body) in enumerate(s["cards"]):
        col = i % 2
        row = i // 2
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        draw_grid_card(slide, x, y, card_w, card_h, icon, title, body, font)
    add_footer(slide, s["footer"], 3, font)


def build_slide4(prs, d):
    """Differentiator: comparison table"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s4"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))
    y = Inches(1.15)
    row_h = Inches(0.5)
    draw_table_row(slide, y, row_h, s["headers"], font, is_header=True)
    y += row_h
    for cells in s["rows"]:
        draw_table_row(slide, y, row_h, cells, font, is_first_col_header=True)
        y += row_h
    add_footer(slide, s["footer"], 4, font)


def build_slide5(prs, d):
    """Partnership models: 3 rows"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s5"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))
    item_w = Inches(8.6)
    item_h = Inches(1.15)
    gap = Inches(0.12)
    start_y = Inches(1.15)
    x = Inches(0.5)
    for i, (badge, color, title, body, ex) in enumerate(s["models"]):
        y = start_y + i * (item_h + gap)
        draw_model_item(slide, x, y, item_w, item_h, badge, color, title, body, ex, font)
    add_footer(slide, s["footer"], 5, font)


def build_slide6(prs, d):
    """Revenue flow + callout"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s6"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))
    box_w = Inches(2.5)
    box_h = Inches(1.05)
    gap = Inches(0.4)
    total = 3 * box_w + 2 * gap
    start_x = (SLIDE_W - total) / 2
    y = Inches(1.2)
    for i, (title, body, bg, tc) in enumerate(s["flows"]):
        x = start_x + i * (box_w + gap)
        draw_flow_box(slide, x, y, box_w, box_h, title, body, bg, tc, font)
        if i < 2:
            arrow_x = x + box_w + Inches(0.05)
            add_textbox(slide, arrow_x, y + Inches(0.3), Inches(0.3), Inches(0.35),
                        "\u2190", font, 18, TEXT_MUTED, align=PP_ALIGN.CENTER)
    # Callout (left accent stripe only, no outer border)
    cx, cy = Inches(0.5), Inches(2.65)
    cw, ch = Inches(8.6), Inches(1.15)
    add_rect(slide, cx, cy, cw, ch, fill=TOKI_BLUE_PALE)
    add_rect(slide, cx, cy, Inches(0.06), ch, fill=TOKI_BLUE)
    add_textbox(slide, cx + Inches(0.25), cy + Inches(0.1), cw - Inches(0.35), Inches(0.28),
                s["callout_title"], font, 11, TEXT_PRIMARY, bold=True)
    add_textbox(slide, cx + Inches(0.25), cy + Inches(0.45), cw - Inches(0.35), ch - Inches(0.55),
                s["callout_body"], font, 10, TEXT_SECONDARY)
    add_footer(slide, s["footer"], 6, font)


def build_slide7(prs, d):
    """Client sectors: 3x2 grid"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s7"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))
    add_textbox(slide, Inches(0.5), Inches(1.1), Inches(8), Inches(0.3),
                s["lead"], font, 11, TEXT_SECONDARY)
    card_w = Inches(2.85)
    card_h = Inches(1.55)
    gap_x = Inches(0.23)
    gap_y = Inches(0.15)
    start_x = Inches(0.5)
    start_y = Inches(1.5)
    for i, (icon, title, body) in enumerate(s["sectors"]):
        col = i % 3
        row = i // 3
        x = start_x + col * (card_w + gap_x)
        y = start_y + row * (card_h + gap_y)
        draw_sector_card(slide, x, y, card_w, card_h, icon, title, body, font)
    add_footer(slide, s["footer"], 7, font)


def build_slide8(prs, d):
    """Founder + independence"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s8"]
    add_action_bar(slide, s["bar"], font)
    add_section_label(slide, s["label"], font, Inches(0.8))

    # ── Profile section ──────────────────────────────────
    # Photo (circular, 1.0" diameter)
    ax, ay, asize = Inches(0.5), Inches(1.12), Inches(1.0)
    photo_path = os.path.join(OUT_DIR, "asset", "IMG_4310-2.jpeg")
    if os.path.exists(photo_path):
        try:
            circ_path = _prepare_circular_photo(photo_path)
            slide.shapes.add_picture(circ_path, ax, ay, asize, asize)
        except Exception:
            _draw_avatar_fallback(slide, ax, ay, asize, d, font)
    else:
        _draw_avatar_fallback(slide, ax, ay, asize, d, font)
    # Name (right of photo)
    name_x = Inches(1.65)
    add_textbox(slide, name_x, Inches(1.18), Inches(7.3), Inches(0.35),
                s["name"], font, 13, TEXT_PRIMARY, bold=True)
    # Bio (right of photo, below name)
    bio_y = 1.58
    add_textbox(slide, name_x, Inches(bio_y), Inches(7.3), Inches(1.2),
                s["bio"], font, 10, TEXT_SECONDARY)

    # ── Tags section (full width, positioned below bio) ──
    # Estimate bio rendered height to avoid overlap
    cpl = 48 if font == FONT_JP else 85  # approx chars per line
    bio_lines = (len(s["bio"]) + cpl - 1) // cpl
    bio_h_est = bio_lines * 0.19  # 10pt * 1.35 line spacing
    tag_start_x = Inches(0.5)
    tag_x = tag_start_x
    tag_y = Inches(max(2.6, bio_y + bio_h_est + 0.35))
    tag_h = Inches(0.44)
    for tag in s["tags"]:
        tw = Inches(len(tag) * 0.075 + 0.45)
        if tag_x + tw > Inches(9.3):
            tag_x = tag_start_x
            tag_y += Inches(0.52)
        add_rect(slide, tag_x, tag_y, tw, tag_h, fill=BG_SECTION, border_color=BORDER)
        add_textbox(slide, tag_x + Inches(0.1), tag_y, tw - Inches(0.2), tag_h,
                    tag, font, 8, TEXT_SECONDARY, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        tag_x += tw + Inches(0.12)

    # ── Independence callout ─────────────────────────────
    cy = tag_y + Inches(0.55)
    cx, cw, ch = Inches(0.5), Inches(8.6), Inches(1.15)
    add_rect(slide, cx, cy, cw, ch, fill=TOKI_BLUE_PALE)
    add_rect(slide, cx, cy, Inches(0.06), ch, fill=TOKI_BLUE)
    add_textbox(slide, cx + Inches(0.25), cy + Inches(0.1), cw - Inches(0.35), Inches(0.28),
                s["ind_title"], font, 10, TEXT_PRIMARY, bold=True)
    add_textbox(slide, cx + Inches(0.25), cy + Inches(0.44), cw - Inches(0.35), ch - Inches(0.52),
                s["ind_body"], font, 10, TEXT_SECONDARY)
    add_footer(slide, s["footer"], 8, font)


def _draw_avatar_fallback(slide, ax, ay, asize, d, font):
    avatar = slide.shapes.add_shape(MSO_SHAPE.OVAL, ax, ay, asize, asize)
    avatar.fill.solid()
    avatar.fill.fore_color.rgb = DARK_BG
    avatar.line.fill.background()
    initials = "佐" if d is CONTENT["ja"] else "TS"
    add_textbox(slide, ax, ay, asize, asize,
                initials, font, 16, WHITE, bold=True, align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)


def build_slide9(prs, d):
    """Next Steps — structured"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s9"]
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=DARK_BG)
    add_textbox(slide, Inches(0.5), Inches(0.15), Inches(4), Inches(0.45),
                s["bar"], font, 13, WHITE, bold=True)
    add_textbox(slide, Inches(0.5), Inches(0.85), Inches(8), Inches(0.5),
                s["title"], font, 20, WHITE, bold=True)
    # 4 numbered steps
    step_y = Inches(1.6)
    step_h = Inches(0.62)
    step_gap = Inches(0.1)
    for i, (num, title, desc) in enumerate(s["steps"]):
        y = step_y + i * (step_h + step_gap)
        circ = slide.shapes.add_shape(MSO_SHAPE.OVAL,
                                       Inches(0.8), y + Inches(0.1),
                                       Inches(0.42), Inches(0.42))
        circ.fill.solid()
        circ.fill.fore_color.rgb = TOKI_BLUE
        circ.line.fill.background()
        add_textbox(slide, Inches(0.8), y + Inches(0.1), Inches(0.42), Inches(0.42),
                    num, font, 10, WHITE, bold=True,
                    align=PP_ALIGN.CENTER, anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(1.4), y + Inches(0.06), Inches(4), Inches(0.3),
                    title, font, 12, WHITE, bold=True)
        add_textbox(slide, Inches(1.4), y + Inches(0.34), Inches(7), Inches(0.26),
                    desc, font, 10, RGBColor(0xBB, 0xBB, 0xCC))
    # Contact
    add_textbox(slide, Inches(0.8), SLIDE_H - Inches(1.0), Inches(8), Inches(0.3),
                s["contact"], font, 10, TEXT_MUTED)
    # Footer
    stripe_y = SLIDE_H - Inches(0.55)
    add_rect(slide, 0, stripe_y - Inches(0.03), SLIDE_W, Pt(0.5), fill=RGBColor(0x33, 0x44, 0x55))
    add_textbox(slide, Inches(0.5), stripe_y, Inches(4), Inches(0.35),
                s["footer_left"], font, 9, TEXT_MUTED)
    add_textbox(slide, Inches(4), stripe_y, Inches(2), Inches(0.35),
                "Confidential", font, 9, TEXT_MUTED, align=PP_ALIGN.CENTER)
    add_textbox(slide, SLIDE_W - Inches(1), stripe_y, Inches(0.5), Inches(0.35),
                "9", font, 9, TEXT_MUTED, bold=True, align=PP_ALIGN.RIGHT)


def build_slide10(prs, d):
    """Disclaimer"""
    slide = add_blank_slide(prs)
    font = d["font"]
    s = d["s10"]
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, fill=BG_PAGE)
    add_textbox(slide, Inches(1), Inches(1.0), Inches(8), Inches(0.45),
                s["title"], font, 14, TEXT_MUTED, bold=True)
    add_rect(slide, Inches(1), Inches(1.5), Inches(8), Pt(0.5), fill=BORDER)
    add_textbox(slide, Inches(1), Inches(1.7), Inches(8), Inches(2.8),
                s["body"], font, 9, TEXT_SECONDARY)
    add_textbox(slide, Inches(1), SLIDE_H - Inches(1.2), Inches(8), Inches(0.3),
                s["copyright"], font, 8, TEXT_MUTED)
    add_footer(slide, s["footer_left"], 10, font)


# ══════════════════════════════════════════════════════════════════════
#  MAIN
# ══════════════════════════════════════════════════════════════════════

def generate(lang):
    d = CONTENT[lang]
    prs = new_prs()
    build_cover(prs, d)
    build_slide2(prs, d)
    build_slide3(prs, d)
    build_slide4(prs, d)
    build_slide5(prs, d)
    build_slide6(prs, d)
    build_slide7(prs, d)
    build_slide8(prs, d)
    build_slide9(prs, d)
    build_slide10(prs, d)

    pptx_path = os.path.join(OUT_DIR, f"{d['filename']}.pptx")
    prs.save(pptx_path)
    _strip_theme_shadows(pptx_path)
    print(f"  PPTX saved: {pptx_path}")
    return pptx_path


def _strip_theme_shadows(pptx_path):
    """Remove default outerShdw / 3D effects from theme1.xml."""
    clean = '<a:effectStyleLst>' + \
            '<a:effectStyle><a:effectLst/></a:effectStyle>' * 3 + \
            '</a:effectStyleLst>'
    buf = io.BytesIO()
    with zipfile.ZipFile(pptx_path, 'r') as zin:
        with zipfile.ZipFile(buf, 'w', zipfile.ZIP_DEFLATED) as zout:
            for item in zin.namelist():
                data = zin.read(item)
                if item == 'ppt/theme/theme1.xml':
                    text = data.decode('utf-8')
                    text = re.sub(r'<a:effectStyleLst>.*?</a:effectStyleLst>',
                                  clean, text, flags=re.DOTALL)
                    data = text.encode('utf-8')
                zout.writestr(item, data)
    with open(pptx_path, 'wb') as f:
        f.write(buf.getvalue())


def convert_to_pdf(pptx_path):
    out_dir = os.path.dirname(pptx_path)
    env = os.environ.copy()
    env["HOME"] = "/tmp"
    result = subprocess.run(
        ["libreoffice", "--headless", "--convert-to", "pdf", "--outdir", out_dir, pptx_path],
        capture_output=True, text=True, timeout=120, env=env
    )
    if result.returncode != 0:
        print(f"  ERROR: {result.stderr}")
        return None
    pdf_path = pptx_path.replace(".pptx", ".pdf")
    print(f"  PDF saved: {pdf_path}")
    return pdf_path


if __name__ == "__main__":
    print("=== TokiStorage Partnership Deck Generator ===\n")

    for lang in ["ja", "en"]:
        print(f"[{lang.upper()}] Generating PPTX...")
        pptx = generate(lang)
        print(f"[{lang.upper()}] Converting to PDF...")
        convert_to_pdf(pptx)

    print("\nDone!")
